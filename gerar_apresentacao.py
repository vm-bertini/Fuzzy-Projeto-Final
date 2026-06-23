"""
Geração da apresentação PPTX — Projeto ANFIS  (FEEC/Unicamp)
Execução: python gerar_apresentacao.py
Saída:    apresentacao_anfis.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Paleta FEEC / Unicamp
# ---------------------------------------------------------------------------
NAVY     = RGBColor(0x1C, 0x3F, 0x6E)   # azul-marinho (texto "feec")
VERDE_D  = RGBColor(0x2E, 0x7D, 0x32)   # verde escuro (logo, canto superior)
VERDE_M  = RGBColor(0x43, 0xA0, 0x47)   # verde médio
VERDE_C  = RGBColor(0xC8, 0xE6, 0xC9)   # verde-menta claro (fundo suave)
TEAL     = RGBColor(0x00, 0x69, 0x5C)   # teal escuro
TEAL_C   = RGBColor(0xB2, 0xDF, 0xDB)   # teal claro
BRANCO   = RGBColor(0xFF, 0xFF, 0xFF)
PRETO    = RGBColor(0x21, 0x21, 0x21)
CINZA    = RGBColor(0x45, 0x5A, 0x64)
CINZA_C  = RGBColor(0xEC, 0xEF, 0xF1)
AMARELO  = RGBColor(0xF9, 0xA8, 0x25)
VERMELHO = RGBColor(0xC6, 0x28, 0x28)

W = Inches(13.33)
H = Inches(7.5)
RESULTS_DIR = os.path.join('src', 'results', '2026_06_23_03_02')


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _rect(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    return s


def _tf(shape, text, size, bold=False, color=BRANCO, align=PP_ALIGN.LEFT, italic=False):
    tf = shape.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tf


def _tb(slide, text, x, y, w, h, size=14, bold=False,
        color=PRETO, align=PP_ALIGN.LEFT, italic=False):
    t = slide.shapes.add_textbox(x, y, w, h)
    tf = t.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return t


def _bullets(slide, items, x, y, w, h, size=14, color=PRETO,
             line_spacing=Pt(5)):
    t  = slide.shapes.add_textbox(x, y, w, h)
    tf = t.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_before = line_spacing
        p.alignment    = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = item
        r.font.size      = Pt(size)
        r.font.color.rgb = color
    return t


def _header(slide, title, sub=None):
    bar = _rect(slide, Inches(0), Inches(0), W, Inches(1.25), NAVY)
    _tf(bar, title, 26, bold=True, color=BRANCO, align=PP_ALIGN.LEFT)
    bar.text_frame.paragraphs[0].runs[0]  # already set
    # indent title
    bar.left   = Inches(0)
    bar.top    = Inches(0)
    bar.width  = W
    bar.height = Inches(1.25)
    # re-position text inside bar with a margin
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.18), W - Inches(0.8), Inches(0.9))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r  = p.add_run(); r.text = title
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = BRANCO
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(12); r2.font.italic = True; r2.font.color.rgb = VERDE_C
    return bar


def _footer(slide, text='Victor M. Bertini  |  FEEC / Unicamp  |  Junho 2026'):
    _rect(slide, Inches(0), H - Inches(0.3), W, Inches(0.3), NAVY)
    _tb(slide, text, Inches(0.3), H - Inches(0.28), W - Inches(0.6), Inches(0.26),
        size=8, color=BRANCO)


def _img(slide, path, x, y, w, h=None):
    if os.path.exists(path):
        if h:
            slide.shapes.add_picture(path, x, y, w, h)
        else:
            slide.shapes.add_picture(path, x, y, w)


def _label(slide, text, x, y, w, h, bg, fg=BRANCO, size=11, bold=True):
    b = _rect(slide, x, y, w, h, bg)
    _tf(b, text, size, bold=bold, color=fg, align=PP_ALIGN.CENTER)
    return b


# ---------------------------------------------------------------------------
# Início
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BL = prs.slide_layouts[6]   # blank


# ==========================================================================
# Slide 1 — Título
# ==========================================================================
s = prs.slides.add_slide(BL)
_rect(s, Inches(0), Inches(0), W, H, NAVY)
_rect(s, Inches(0), Inches(2.55), W, Inches(2.6), VERDE_D)

_tb(s, 'Sistemas Neuro-Fuzzy para Classificação de Risco Cardíaco',
    Inches(0.6), Inches(2.65), W - Inches(1.2), Inches(1.3),
    size=32, bold=True, color=BRANCO, align=PP_ALIGN.LEFT)

_tb(s, 'Comparação entre ANFIS e MLP: Acurácia vs Interpretabilidade',
    Inches(0.6), Inches(3.9), W - Inches(1.2), Inches(0.65),
    size=17, color=VERDE_C, align=PP_ALIGN.LEFT, italic=True)

_tb(s, 'Victor M. Bertini', Inches(0.6), Inches(5.0),
    Inches(6), Inches(0.5), size=15, bold=True, color=BRANCO)
_tb(s, 'FEEC / Unicamp  —  Lógica Fuzzy  —  Junho 2026',
    Inches(0.6), Inches(5.45), Inches(8), Inches(0.4), size=13, color=VERDE_C)


# ==========================================================================
# Slide 2 — Motivação
# ==========================================================================
s = prs.slides.add_slide(BL)
_header(s, 'Motivação: O Custo da Caixa-Preta')
_footer(s)

_bullets(s, [
    '▸  Modelos de ML alcançam alta acurácia em diagnóstico cardíaco',
    '▸  Médicos precisam ENTENDER a decisão — não apenas confiar nela',
    '▸  "Por que este paciente foi classificado como alto risco?"',
    '▸  Regulações e protocolos clínicos exigem rastreabilidade das decisões',
    '▸  Redes neurais tradicionais (MLP) não fornecem nenhuma explicação',
    '',
    '→  Pergunta central:',
    '    É possível ter desempenho comparável ao MLP',
    '    COM interpretabilidade linguística?',
], Inches(0.7), Inches(1.4), W - Inches(1.4), H - Inches(1.8), size=16, color=PRETO)


# ==========================================================================
# Slide 3 — Lacuna
# ==========================================================================
s = prs.slides.add_slide(BL)
_header(s, 'A Lacuna: Acurácia vs Interpretabilidade')
_footer(s)

# Caixa MLP
b1 = _rect(s, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.6), CINZA_C)
_tb(s, 'MLP  —  Caixa-Preta', Inches(0.6), Inches(1.5),
    Inches(5.4), Inches(0.5), size=15, bold=True, color=NAVY)
_bullets(s, [
    '✓  Alta acurácia  (AUC ≈ 0,91)',
    '✓  Treinamento rápido  (~5 s)',
    '✗  Zero interpretabilidade',
    '✗  Nenhuma regra auditável',
    '✗  Inaceitável em contextos regulados',
], Inches(0.6), Inches(2.05), Inches(5.4), Inches(4.5), size=14, color=PRETO)

# Seta
_tb(s, '←  vs  →', Inches(6.27), Inches(4.0), Inches(0.85), Inches(0.5),
    size=18, bold=True, color=CINZA, align=PP_ALIGN.CENTER)

# Caixa ANFIS
b2 = _rect(s, Inches(7.15), Inches(1.4), Inches(5.8), Inches(5.6), VERDE_C)
_tb(s, 'ANFIS  —  Neuro-Fuzzy', Inches(7.35), Inches(1.5),
    Inches(5.4), Inches(0.5), size=15, bold=True, color=VERDE_D)
_bullets(s, [
    '✓  Acurácia comparável  (AUC ≈ 0,90)',
    '✓  Regras linguísticas auditáveis',
    '✓  Parâmetros com significado físico',
    '✓  13× menos parâmetros que MLP',
    '~  Treinamento ligeiramente mais lento',
], Inches(7.35), Inches(2.05), Inches(5.4), Inches(4.5), size=14, color=PRETO)


# ==========================================================================
# Slide 4 — Roteiro
# ==========================================================================
s = prs.slides.add_slide(BL)
_header(s, 'Estrutura da Apresentação')
_footer(s)

_bullets(s, [
    '1.   Framework Teórico — arquitetura ANFIS camada a camada',
    '2.   Dataset e Configuração Experimental',
    '3.   Pipeline K-Means — geração do pool de regras',
    '4.   Métricas de Avaliação utilizadas',
    '5.   MFs de Saída — conversão TS → rótulo linguístico',
    '6.   Resultados: desempenho, MFs aprendidas e regras extraídas',
    '7.   Análise Comparativa ANFIS × MLP',
    '8.   Conclusões e Trabalho Futuro',
], Inches(1.2), Inches(1.45), W - Inches(2.0), H - Inches(1.85), size=17, color=PRETO)


# ==========================================================================
# Slide 5 — Framework Teórico
# ==========================================================================
s = prs.slides.add_slide(BL)
_header(s, 'Arquitetura ANFIS — Camada a Camada')
_footer(s)

Y0 = Inches(1.35)
bh = Inches(0.55)
bw = Inches(2.55)

# --- Pipeline de caixas no topo ---
boxes = [
    ('Entrada\n[B, 11]',             Inches(0.2),  NAVY),
    ('MixedFuzzyLayer\n[B, 11, M]',  Inches(2.85), VERDE_D),
    ('ClusterRuleLayer\n[B, R]',      Inches(5.5),  TEAL),
    ('DefuzzyLayer\n[B, 1]',         Inches(8.15), VERDE_D),
    ('Sigmoid → P(d)',                Inches(10.8), NAVY),
]
for lbl, bx, col in boxes:
    b = _rect(s, bx, Y0, bw, bh, col)
    _tf(b, lbl, 9, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    if bx < Inches(10.8):
        _tb(s, '→', bx + bw, Y0 + Inches(0.08), Inches(0.28), Inches(0.38),
            size=14, bold=True, color=CINZA, align=PP_ALIGN.CENTER)

# --- Descrição por camada (3 colunas) ---
col_x = [Inches(0.2), Inches(4.5), Inches(8.85)]
col_w  = Inches(4.1)

# Coluna 1: MixedFuzzyLayer
_tb(s, 'MixedFuzzyLayer', col_x[0], Inches(2.1), col_w, Inches(0.38),
    size=13, bold=True, color=VERDE_D)
_bullets(s, [
    'Numéricas (5): MFs Gaussianas treináveis',
    '  μ(x) = exp(−½·((x−c)/σ)²)',
    '  c, σ ∈ ℝ  ajustados por backprop',
    '',
    'Categóricas (6): one-hot suavizado',
    '  out[cat certa] = 0,95',
    '  out[demais]   = 0,05/(n−1)',
    '',
    'Exemplos categóricos:',
    '  ChestPainType: ASY→3, ATA→0, NAP→1',
    '  ExerciseAngina: N→0, Y→1',
    '  ST_Slope: Down→0, Flat→1, Up→2',
], col_x[0], Inches(2.5), col_w, Inches(4.5), size=11, color=PRETO)

# Coluna 2: ClusterRuleLayer
_tb(s, 'ClusterRuleLayer', col_x[1], Inches(2.1), col_w, Inches(0.38),
    size=13, bold=True, color=TEAL)
_bullets(s, [
    'T-norma produto (diferenciável):',
    '',
    '  w_k = ∏ᵢ μᵢ,combo[k,i](xᵢ)',
    '',
    'combo[k,i] = índice da MF da regra k',
    'na feature i (gerado pelo K-Means)',
    '',
    'Normalização:',
    '  w̃_k = w_k / (Σⱼ wⱼ + ε)',
    '',
    'Sem parâmetros treináveis.',
    'rule_idx salvo como register_buffer',
    '(move para GPU automaticamente)',
], col_x[1], Inches(2.5), col_w, Inches(4.5), size=11, color=PRETO)

# Coluna 3: DefuzzyLayer + Loss
_tb(s, 'DefuzzyLayer  +  Loss', col_x[2], Inches(2.1), col_w, Inches(0.38),
    size=13, bold=True, color=VERDE_D)
_bullets(s, [
    'Takagi-Sugeno (ordem 0):',
    '  ŷ = Σ_k w̃_k · c_k',
    '  c_k ∈ nn.Linear(R→1) — consequentes',
    '',
    'Loss = BCE + λ·penalidade',
    '',
    '  BCE(σ(ŷ), y)  — erro de classificação',
    '',
    '  λ·Σᵢ<ⱼ relu(σᵢ+σⱼ−|cᵢ−cⱼ|)',
    '  → penaliza MFs sobrepostas:',
    '  se |cᵢ−cⱼ| < σᵢ+σⱼ, MF j cabe',
    '  dentro de MF i → partição incoerente',
    '  λ=0,1 mantém como regularizador',
], col_x[2], Inches(2.5), col_w, Inches(4.5), size=11, color=PRETO)


# ==========================================================================
# Slide 6 — Dataset + enxerto
# ==========================================================================
s = prs.slides.add_slide(BL)
_header(s, 'Dataset e Configuração Experimental')
_footer(s)

# Coluna esquerda — info
_tb(s, 'Heart Failure Prediction Dataset',
    Inches(0.5), Inches(1.4), Inches(5.5), Inches(0.45),
    size=13, bold=True, color=NAVY)
_bullets(s, [
    '918 pacientes adultos',
    'Variável alvo: HeartDisease (0/1)',
    'Prevalência: ~55% positivos',
    'Split: 80% treino / 20% teste  (≈734 / 184)',
    '',
    'Features numéricas (5):',
    '  Age, RestingBP, Cholesterol,',
    '  MaxHR, Oldpeak',
    '',
    'Features categóricas (6):',
    '  Sex, ChestPainType, FastingBS,',
    '  RestingECG, ExerciseAngina, ST_Slope',
    '',
    'Configuração:',
    '  NUM_MFS=3  |  EPOCHS=200  |  lr=1e-3',
    '  ANFIS pool: 200 clusters K-Means',
    '  MLP: camadas [64, 32]',
], Inches(0.5), Inches(1.9), Inches(5.5), H - Inches(2.5), size=12, color=PRETO)

# Coluna direita — snippet da tabela
_tb(s, 'Exemplo de registros no dataset:',
    Inches(6.2), Inches(1.4), Inches(6.7), Inches(0.4),
    size=12, bold=True, color=NAVY)

# Montar tabela manual (Age, Sex, ChestPain, MaxHR, ExAngina, OldPeak, ST, HD)
headers = ['Age', 'Sex', 'ChestPain', 'MaxHR', 'ExAngina', 'Oldpeak', 'ST_Slope', 'HD']
rows_data = [
    ['40', 'M', 'ATA', '172', 'N', '0.0', 'Up',   '0'],
    ['49', 'F', 'NAP', '156', 'N', '1.0', 'Flat', '1'],
    ['37', 'M', 'ATA', ' 98', 'N', '0.0', 'Up',   '0'],
    ['48', 'F', 'ASY', '108', 'Y', '1.5', 'Flat', '1'],
]

cell_w = Inches(0.78)
cell_h = Inches(0.42)
tx0    = Inches(6.2)
ty0    = Inches(1.85)

for ci, hdr in enumerate(headers):
    b = _rect(s, tx0 + ci*cell_w, ty0, cell_w - Inches(0.02), cell_h, NAVY)
    _tf(b, hdr, 9, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows_data):
    bg = VERDE_C if ri % 2 == 0 else BRANCO
    for ci, val in enumerate(row):
        b = _rect(s, tx0 + ci*cell_w, ty0 + (ri+1)*cell_h,
                  cell_w - Inches(0.02), cell_h - Inches(0.02), bg)
        _tf(b, val, 9, color=PRETO, align=PP_ALIGN.CENTER)

# Nota abaixo da tabela
_tb(s, '* HD = HeartDisease (variável alvo)',
    Inches(6.2), ty0 + Inches(2.1), Inches(6.7), Inches(0.35),
    size=10, color=CINZA, italic=True)

# Pré-processamento
_tb(s, 'Pré-processamento:',
    Inches(6.2), Inches(4.45), Inches(6.7), Inches(0.4),
    size=12, bold=True, color=NAVY)
_bullets(s, [
    '▸  PartialScaler: z-score só nas features numéricas',
    '▸  Categóricas mantêm código inteiro (pd.Categorical.codes)',
    '▸  Cholesterol/RestingBP=0 → substituído pela mediana',
], Inches(6.2), Inches(4.85), Inches(6.7), Inches(1.8), size=12, color=PRETO)


# ==========================================================================
# Slide 7 — K-Means pipeline
# ==========================================================================
s = prs.slides.add_slide(BL)
_header(s, 'Geração do Pool de Regras via K-Means')
_footer(s)

# Caixa de aviso no topo
av = _rect(s, Inches(0.4), Inches(1.35), W - Inches(0.8), Inches(0.7), AMARELO)
_tf(av,
    'Importante: a fuzzificação percentílica abaixo é um ARTIFÍCIO de mapeamento '
    '— não são as MFs do ANFIS. Servem exclusivamente para selecionar quais combinações '
    'de antecedentes existem no pool de regras.',
    10, bold=False, color=PRETO, align=PP_ALIGN.LEFT)

_bullets(s, [
    '1.  Picos de MF por percentil  (generalização para facilitar o mapeamento):',
    '      pcts = [25%, 50%, 75%]  →  peaks[fi] = percentis da feature fi no treino',
    '      Cada cluster center é mapeado para baixo / médio / alto por argmax da MF triangular',
    '      → Artifício: converte vetor contínuo em índice discreto de MF',
    '',
    '2.  K-Means (200 clusters) particiona o espaço de entrada em regiões densas',
    '      Sem supervisão — apenas agrupa por proximidade no espaço das features',
    '',
    '3.  Mapeamento dominante por cluster:',
    '      Numérica:   argmax μ_triangular(centro_k, pico_j)  →  índice 0,1,2',
    '      Categórica: round(centro_k)                        →  código inteiro',
    '',
    '4.  Deduplicação: clusters com mesmo vetor de MFs → uma regra única',
    '      200 clusters → ~80–120 combos únicos',
    '',
    '5.  Centros Gaussianos do ANFIS inicializados nos percentis dos centros K-Means',
    '      Treino end-to-end refina c, σ e c_k independentemente deste mapeamento',
], Inches(0.5), Inches(2.1), W - Inches(1.0), H - Inches(2.65), size=12, color=PRETO)


# ==========================================================================
# Slide 8 — Métricas de benchmark
# ==========================================================================
s = prs.slides.add_slide(BL)
_header(s, 'Métricas de Avaliação Utilizadas')
_footer(s)

metrics = [
    ('AUC-ROC',
     'Área sob a curva ROC (Receiver Operating Characteristic).',
     'Mede a capacidade do modelo de separar as classes independentemente do limiar.',
     'AUC = 1,0 → separação perfeita  |  AUC = 0,5 → equivale a chute aleatório',
     VERDE_D),
    ('Acurácia',
     'Proporção de predições corretas sobre o total de amostras.',
     '  Acurácia = (VP + VN) / N',
     'Sensível ao desbalanceamento de classes — complementada pelo F1.',
     TEAL),
    ('F1-Score',
     'Média harmônica entre Precisão e Revocação.',
     '  F1 = 2 · (Precisão · Revocação) / (Precisão + Revocação)',
     'Preferível à acurácia quando as classes estão desbalanceadas.',
     NAVY),
]

for i, (name, l1, l2, l3, col) in enumerate(metrics):
    bx = Inches(0.4) + i * Inches(4.25)
    b  = _rect(s, bx, Inches(1.4), Inches(4.1), Inches(0.55), col)
    _tf(b, name, 16, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    _bullets(s, [l1, l2, l3],
             bx, Inches(2.0), Inches(4.1), Inches(4.6),
             size=12, color=PRETO)

_tb(s, 'Limiar de decisão em todas as métricas de classificação: 0,5',
    Inches(0.5), Inches(6.85), W - Inches(1.0), Inches(0.35),
    size=11, italic=True, color=CINZA)


# ==========================================================================
# Slide 9 — MFs de saída hardcoded
# ==========================================================================
s = prs.slides.add_slide(BL)
_header(s, 'MFs de Saída — Conversão do Resultado TS em Rótulo Linguístico')
_footer(s)

# Explicação esquerda
_tb(s, 'O problema:', Inches(0.5), Inches(1.4), Inches(6.0), Inches(0.4),
    size=13, bold=True, color=NAVY)
_bullets(s, [
    'A camada DefuzzyLayer produz um valor contínuo ŷ ∈ ℝ.',
    'Após sigmoid: P(doença) ∈ [0, 1].',
    'Regras Takagi-Sugeno não têm consequente linguístico —',
    '  apenas um número (c_k aprendido).',
    '',
    'Para transcrever regras do tipo',
    '  "SE ... ENTÃO risco = ALTO"',
    'precisamos converter ŷ em um rótulo discreto.',
], Inches(0.5), Inches(1.85), Inches(6.0), Inches(3.0), size=12, color=PRETO)

_tb(s, 'A solução (análise manual):', Inches(0.5), Inches(4.9), Inches(6.0), Inches(0.4),
    size=13, bold=True, color=NAVY)
_bullets(s, [
    '5 MFs triangulares uniformemente espaçadas em [0,1]:',
    '  centros:  c = [1/6, 2/6, 3/6, 4/6, 5/6]',
    '  meia-largura:  hw = 1/6  (partição exata, sem sobreposição)',
    '',
    '  μⱼ(ŷ) = max(0,  1 − |ŷ − cⱼ| / hw)',
    '',
    '  rótulo(ŷ) = argmaxⱼ  μⱼ(ŷ)',
    '',
    'Ex.: ŷ = 0,78  →  μ₄(alto) = 0,68  →  risco = alto',
], Inches(0.5), Inches(5.35), Inches(6.0), Inches(1.8), size=12, color=PRETO)

# Imagem das MFs de saída à direita
_img(s, os.path.join(RESULTS_DIR, 'mfs_anfis.png'),
     Inches(6.6), Inches(1.35), Inches(6.4), Inches(5.8))

_tb(s, '↑  painel "Output — risco" (canto inferior direito)',
    Inches(6.6), Inches(7.1), Inches(6.4), Inches(0.3),
    size=10, italic=True, color=CINZA)


# ==========================================================================
# Slide 10 — Resultados: métricas
# ==========================================================================
s = prs.slides.add_slide(BL)
_header(s, 'Resultados: Métricas de Desempenho')
_footer(s)

rows = [
    ('Modelo',  'Parâmetros', 'Treino (s)', 'Acurácia', 'F1-Score', 'AUC-ROC'),
    ('ANFIS',   '228',        '13,8 s',     '83,2 %',   '84,7 %',  '0,903'),
    ('MLP',     '2.881',      '4,6 s',      '83,7 %',   '85,0 %',  '0,910'),
]
cx = [Inches(0.4), Inches(2.6), Inches(4.5), Inches(6.4), Inches(8.5), Inches(10.6)]
cw = [Inches(2.1), Inches(1.85), Inches(1.85), Inches(2.0), Inches(2.0), Inches(2.0)]
ry = [Inches(1.55), Inches(2.4), Inches(3.25)]
rh = Inches(0.78)

for ri, row in enumerate(rows):
    bg = NAVY if ri == 0 else (VERDE_C if ri == 1 else CINZA_C)
    fg = BRANCO if ri == 0 else PRETO
    for ci, (cell, x, w) in enumerate(zip(row, cx, cw)):
        b = _rect(s, x, ry[ri], w - Inches(0.04), rh - Inches(0.04), bg)
        _tf(b, cell, 14, bold=(ri == 0), color=fg, align=PP_ALIGN.CENTER)

_bullets(s, [
    '▸  ANFIS: 228 parâmetros  vs  MLP: 2.881  →  13× mais eficiente',
    '▸  Diferença de AUC: apenas Δ = 0,007  (estatisticamente negligenciável)',
    '▸  ANFIS perde levemente em velocidade de treino (~3× mais lento)',
    '▸  ANFIS compensa com regras auditáveis — MLP não oferece nenhuma explicação',
], Inches(0.5), Inches(4.25), W - Inches(1.0), Inches(2.7), size=14, color=PRETO)


# ==========================================================================
# Slide 11 — MFs Aprendidas
# ==========================================================================
s = prs.slides.add_slide(BL)
_header(s, 'Funções de Pertinência Aprendidas pelo ANFIS (pós-treinamento)')
_footer(s)
_img(s, os.path.join(RESULTS_DIR, 'mfs_anfis.png'),
     Inches(0.2), Inches(1.3), W - Inches(0.4), Inches(5.85))


# ==========================================================================
# Slide 12 — Regras linguísticas (todas as 10)
# ==========================================================================
s = prs.slides.add_slide(BL)
_header(s, 'Regras Linguísticas Extraídas — Top-10 por Força de Disparo Média')
_footer(s)

regras = [
    ('R1  ALTO',
     'Age=idoso, ChestPain=ASY, MaxHR=baixa, ExAngina=Y, Oldpeak=elevado, ST=Flat'),
    ('R2  BAIXO',
     'Age=jovem, ChestPain=ATA, MaxHR=alta,  ExAngina=N, Oldpeak=normal,  ST=Up'),
    ('R3  ALTO',
     'Age=idoso, ChestPain=ASY, Chol=desejável, MaxHR=baixa, ExAngina=Y, Oldpeak=elevado, ST=Flat'),
    ('R4  ALTO',
     'Age=jovem, ChestPain=ASY, MaxHR=baixa, ExAngina=Y, Oldpeak=elevado, ST=Flat'),
    ('R5  BAIXO',
     'Age=jovem, ChestPain=NAP, MaxHR=alta,  ExAngina=N, Oldpeak=normal,  ST=Up'),
    ('R6  MÉD_ALTO',
     'Age=idoso, ChestPain=ASY, RestingBP=ótimo, Chol=alto, MaxHR=baixa, ExAngina=Y, ST=Flat'),
    ('R7  BAIXO',
     'Age=jovem, ChestPain=ATA, RestingBP=ótimo, MaxHR=alta, ExAngina=N, Oldpeak=normal, ST=Up'),
    ('R8  ALTO',
     'Age=idoso, ChestPain=NAP, MaxHR=baixa, ExAngina=Y, Oldpeak=elevado, ST=Flat'),
    ('R9  BAIXO_M',
     'Age=jovem, ChestPain=ASY, RestingBP=ótimo, Chol=alto, MaxHR=alta, ExAngina=N, ST=Up'),
    ('R10 BAIXO',
     'Age=jovem, ChestPain=ATA, RestingBP=ótimo, MaxHR=baixa, ExAngina=N, Oldpeak=normal, ST=Up'),
]

rh2  = Inches(0.46)
ty   = Inches(1.35)
col_label_w = Inches(1.5)
col_rule_w  = W - col_label_w - Inches(0.7)

for i, (lbl, rule) in enumerate(regras):
    risk = lbl.split()[-1]
    bg = VERDE_D if 'ALTO' in risk and 'BAIXO' not in risk else \
         TEAL    if 'MÉD' in risk else \
         NAVY
    b = _rect(s, Inches(0.35), ty + i*rh2, col_label_w, rh2 - Inches(0.03), bg)
    _tf(b, lbl, 9, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    _tb(s, rule, col_label_w + Inches(0.45), ty + i*rh2 + Inches(0.05),
        col_rule_w, rh2, size=10, color=PRETO)


# ==========================================================================
# Slide 13 — Comparação de risco
# ==========================================================================
s = prs.slides.add_slide(BL)
_header(s, 'Comparação de Risco: ANFIS vs MLP por Paciente (Conjunto de Teste)')
_footer(s)
_img(s, os.path.join(RESULTS_DIR, 'risk_comparison.png'),
     Inches(0.2), Inches(1.3), W - Inches(0.4), Inches(5.85))


# ==========================================================================
# Slide 14 — Análise crítica
# ==========================================================================
s = prs.slides.add_slide(BL)
_header(s, 'Análise: ANFIS como Alternativa Interpretável ao MLP')
_footer(s)

hl = _rect(s, Inches(0.4), Inches(1.4), W - Inches(0.8), Inches(0.85), VERDE_C)
_tf(hl,
    'ANFIS perde apenas 0,007 de AUC em relação ao MLP, com 13× menos parâmetros '
    'e regras linguísticas auditáveis por especialistas.',
    14, bold=True, color=VERDE_D, align=PP_ALIGN.LEFT)

_bullets(s, [
    'Eficiência paramétrica:',
    '  228 vs 2.881 parâmetros — centros e variâncias das MFs têm interpretação direta',
    '  Cada c_k aprendido indica a "opinião" daquela regra sobre a saída',
    '',
    'Interpretabilidade genuína:',
    '  ST_Slope=Flat + ExerciseAngina=Y + Oldpeak=elevado → padrão recorrente de alto risco',
    '  Compatível com literatura médica sobre isquemia cardíaca',
    '',
    'Tradeoff velocidade:',
    '  ANFIS: 13,8 s (200 épocas)  vs  MLP: 4,6 s',
    '  Diferença aceitável em produção offline / ambiente clínico',
], Inches(0.5), Inches(2.35), W - Inches(1.0), Inches(4.5), size=13, color=PRETO)


# ==========================================================================
# Slide 15 — Pontos de atenção
# ==========================================================================
s = prs.slides.add_slide(BL)
_header(s, 'Pontos de Atenção e Limitações')
_footer(s)

_bullets(s, [
    '1.  Explosão de regras',
    '      n_regras ∝ num_mfs^n_inputs — inviável sem seleção/deduplicação',
    '      Solução aplicada: K-Means como filtro + deduplicação por combo de MFs',
    '',
    '2.  Features categóricas em sistemas fuzzy',
    '      Não existe grau de pertinência natural entre categorias discretas',
    '      MixedFuzzyLayer usa one-hot suavizado como aproximação prática',
    '      Formalmente: μ(cat correta) = 0,95;  μ(demais) = 0,05 / (n−1)',
    '',
    '3.  Sensibilidade à inicialização',
    '      K-Means determinístico (semente fixa) — pool de regras pode variar',
    '      Penalidade de separação mitiga colapso de MFs durante treinamento',
    '',
    '4.  Testado em um único dataset (n=918)',
    '      Generalização para outros domínios requer validação adicional',
], Inches(0.5), Inches(1.4), W - Inches(1.0), H - Inches(1.85), size=13, color=PRETO)


# ==========================================================================
# Slide 16 — Conclusões
# ==========================================================================
s = prs.slides.add_slide(BL)
_header(s, 'Conclusões')
_footer(s)

_bullets(s, [
    '✅  ANFIS atinge AUC 0,903 — comparável ao MLP (0,910) no dataset Heart Failure',
    '✅  228 parâmetros treináveis vs 2.881 do MLP  (13× mais eficiente)',
    '✅  Regras linguísticas auditáveis: centros Gaussianos revelam limiares clínicos',
    '✅  Pipeline reproduzível: K-Means → inicialização → treino end-to-end',
    '✅  Penalidade de separação garante MFs linguisticamente coerentes',
    '',
    '→  ANFIS é uma alternativa viável quando interpretabilidade é requisito,',
    '    sem sacrifício significativo de desempenho preditivo.',
], Inches(0.6), Inches(1.45), W - Inches(1.2), H - Inches(1.9), size=16, color=PRETO)


# ==========================================================================
# Slide 17 — Trabalho Futuro
# ==========================================================================
s = prs.slides.add_slide(BL)
_header(s, 'Limitações e Trabalho Futuro')
_footer(s)

# duas colunas
_tb(s, 'Limitações atuais:', Inches(0.5), Inches(1.4), Inches(6.0), Inches(0.4),
    size=13, bold=True, color=NAVY)
_bullets(s, [
    '• Testado em único dataset (n=918)',
    '• Categóricas sem MFs formalmente definidas',
    '• Pool de regras depende da semente K-Means',
    '• Treinamento ~3× mais lento que MLP',
], Inches(0.5), Inches(1.85), Inches(6.0), Inches(2.5), size=13, color=PRETO)

_tb(s, 'Próximos passos:', Inches(0.5), Inches(4.45), Inches(6.0), Inches(0.4),
    size=13, bold=True, color=NAVY)
_bullets(s, [
    '• Validar em datasets maiores e multi-classe',
    '• Incorporar conhecimento de especialistas',
    '  na inicialização das MFs Gaussianas',
    '• Comparar com XGBoost + SHAP',
    '• Interface de extração em linguagem natural',
], Inches(0.5), Inches(4.9), Inches(6.0), Inches(2.2), size=13, color=PRETO)

_img(s, os.path.join(RESULTS_DIR, 'confusion_matrix.png'),
     Inches(6.5), Inches(1.35), Inches(6.5), Inches(5.8))


# ==========================================================================
# Slide 18 — Obrigado
# ==========================================================================
s = prs.slides.add_slide(BL)
_rect(s, Inches(0), Inches(0), W, H, NAVY)
_rect(s, Inches(0), Inches(2.7), W, Inches(2.2), VERDE_D)

_tb(s, 'Obrigado!',
    Inches(0.6), Inches(2.8), W - Inches(1.2), Inches(1.1),
    size=46, bold=True, color=BRANCO, align=PP_ALIGN.LEFT)

_tb(s, 'victor.m.bertini@gmail.com',
    Inches(0.6), Inches(4.1), W, Inches(0.5),
    size=16, color=VERDE_C)

_tb(s, 'github.com/vm-bertini/Fuzzy-Projeto-Final',
    Inches(0.6), Inches(4.65), W, Inches(0.5),
    size=14, color=VERDE_C)

_tb(s, '"ANFIS: interpretabilidade sem sacrificar performance"',
    Inches(0.6), Inches(5.5), W, Inches(0.5),
    size=14, italic=True, color=BRANCO)


# ==========================================================================
# Apêndice A — Heatmap de Regras
# ==========================================================================
s = prs.slides.add_slide(BL)
_header(s, 'Apêndice A — Heatmap das Regras ANFIS (Top-10)')
_footer(s)
_img(s, os.path.join(RESULTS_DIR, 'rule_comparison.png'),
     Inches(0.2), Inches(1.3), W - Inches(0.4), Inches(5.85))


# ==========================================================================
# Apêndice B — Partição K-Means
# ==========================================================================
s = prs.slides.add_slide(BL)
_header(s, 'Apêndice B — Partição K-Means (pré-ANFIS)')
_footer(s)
_img(s, os.path.join(RESULTS_DIR, 'mfs_kmeans.png'),
     Inches(0.2), Inches(1.3), W - Inches(0.4), Inches(5.85))


# ---------------------------------------------------------------------------
# Salvar
# ---------------------------------------------------------------------------
out = 'apresentacao_anfis_v2.pptx'
prs.save(out)
print(f'Apresentação salva: {out}  ({len(prs.slides)} slides)')
